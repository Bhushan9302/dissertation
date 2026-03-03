"""
generate_presentation.py
------------------------
Generates a professional PowerPoint presentation (.pptx) for the MSc
dissertation titled "Mapping and Modelling the Spatial Diffusion of AI
in the UK".

Usage:
    python generate_presentation.py

Output:
    dissertation_presentation.pptx  (in the repository root)

Requirements:
    pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x00, 0x33, 0x66)   # primary header / title colour
MID_BLUE    = RGBColor(0x00, 0x5F, 0xAD)   # secondary accent
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)   # subtle background tint
ACCENT_GOLD = RGBColor(0xF0, 0xA5, 0x00)   # highlight / emphasis
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xC0, 0x39, 0x2B)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def set_bg_color(slide, color: RGBColor):
    """Fill a slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor,
             line_color: RGBColor = None):
    """Add a filled rectangle shape to a slide."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, italic=False,
                 color: RGBColor = DARK_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box with a single paragraph to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a dark-blue title bar at the top of a content slide."""
    bar = add_rect(slide, Inches(0), Inches(0),
                   SLIDE_W, Inches(1.3), DARK_BLUE)

    # Title text inside bar
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(14)
        run2.font.italic = True
        run2.font.color.rgb = LIGHT_BLUE


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=DARK_GRAY, indent_level=0):
    """Add a bullet list text box; items can be (text, level) tuples or strings."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, indent_level

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = level
        bullet_char = "▪" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

    return txBox


def add_footer(slide, page_num: int, total: int = 18):
    """Add a thin footer bar with slide number."""
    add_rect(slide, Inches(0), Inches(7.2),
             SLIDE_W, Inches(0.3), DARK_BLUE)
    add_text_box(
        slide,
        f"Mapping and Modelling the Spatial Diffusion of AI in the UK  |  "
        f"MSc Dissertation  |  Bhushan  |  Slide {page_num} of {total}",
        Inches(0.3), Inches(7.2), Inches(12.7), Inches(0.3),
        font_size=9, color=LIGHT_BLUE, align=PP_ALIGN.LEFT
    )


# ---------------------------------------------------------------------------
# Individual slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs: Presentation):
    """Slide 1 – Title Slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg_color(slide, DARK_BLUE)

    # Gold accent bar (top)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ACCENT_GOLD)
    # Gold accent bar (bottom)
    add_rect(slide, Inches(0), Inches(7.38), SLIDE_W, Inches(0.12), ACCENT_GOLD)

    # Central white panel
    add_rect(slide, Inches(1.0), Inches(1.2),
             Inches(11.33), Inches(5.1), WHITE)

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1.3), Inches(1.5), Inches(10.73), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Mapping and Modelling the"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Spatial Diffusion of AI in the UK"
    run2.font.size = Pt(36)
    run2.font.bold = True
    run2.font.color.rgb = DARK_BLUE

    # Gold divider line
    add_rect(slide, Inches(3.0), Inches(3.45), Inches(7.33), Inches(0.04),
             ACCENT_GOLD)

    # Subtitle
    add_text_box(
        slide, "MSc Dissertation Project",
        Inches(1.3), Inches(3.6), Inches(10.73), Inches(0.5),
        font_size=22, bold=False, italic=True,
        color=MID_BLUE, align=PP_ALIGN.CENTER
    )

    # Author & date
    year = datetime.date.today().year
    add_text_box(
        slide, "Bhushan",
        Inches(1.3), Inches(4.3), Inches(10.73), Inches(0.4),
        font_size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, f"MSc in Applied Data Science  |  {year}",
        Inches(1.3), Inches(4.75), Inches(10.73), Inches(0.35),
        font_size=14, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER
    )

    # Keywords ribbon
    kw_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(5.55), Inches(11.33), Inches(0.55))
    tf_kw = kw_box.text_frame
    p_kw = tf_kw.paragraphs[0]
    p_kw.alignment = PP_ALIGN.CENTER
    run_kw = p_kw.add_run()
    run_kw.text = (
        "Machine Learning   ·   Spatial Analysis   ·   "
        "Location Quotients   ·   LDA Topic Modelling   ·   UK AI Economy"
    )
    run_kw.font.size = Pt(12)
    run_kw.font.italic = True
    run_kw.font.color.rgb = MID_GRAY


def slide_02_toc(prs: Presentation):
    """Slide 2 – Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Table of Contents")
    add_footer(slide, 2)

    sections = [
        ("1.  Research Context & Motivation", 0),
        ("2.  Research Questions & Objectives", 0),
        ("3.  Literature Review – Theoretical Frameworks", 0),
        ("4.  Data & Data Collection", 0),
        ("5.  Methodology Overview", 0),
        ("6.  Entity Resolution & HQ Proxy Modelling", 0),
        ("7.  AI Classification Model (95.32% accuracy)", 0),
        ("8.  Topic Modelling – LDA Results", 0),
        ("9.  Spatial Analysis Methodology", 0),
        ("10. Key Findings – Spatial Patterns", 0),
        ("11. Key Findings – Sectoral Analysis", 0),
        ("12. Key Findings – Diffusion Patterns", 0),
        ("13. Discussion & Implications", 0),
        ("14. Limitations & Future Work", 0),
        ("15. Conclusion", 0),
    ]

    col1 = sections[:8]
    col2 = sections[8:]

    # Left column
    add_bullet_list(
        slide, [t for t, _ in col1],
        Inches(0.4), Inches(1.5), Inches(6.1), Inches(5.7),
        font_size=14, color=DARK_GRAY
    )
    # Right column
    add_bullet_list(
        slide, [t for t, _ in col2],
        Inches(6.7), Inches(1.5), Inches(6.1), Inches(5.7),
        font_size=14, color=DARK_GRAY
    )


def slide_03_context(prs: Presentation):
    """Slide 3 – Research Context & Motivation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Research Context & Motivation",
                  "Why study the spatial diffusion of AI in the UK?")
    add_footer(slide, 3)

    bullets = [
        "Artificial Intelligence has emerged as a transformative, "
        "general-purpose technology reshaping entire industries.",
        "The UK is one of the leading AI economies globally, yet "
        "geographic patterns of firm-level AI adoption remain poorly understood.",
        "Existing research focuses on AI investment flows or national aggregates; "
        "granular spatial analysis is lacking.",
        "Understanding where AI firms cluster—and why—has direct implications "
        "for regional policy, infrastructure investment, and digital inclusion.",
        "This dissertation fills the gap by mapping AI diffusion at the "
        "UK Outcode (Postal District) level using firm-level web data.",
    ]

    add_bullet_list(
        slide, bullets,
        Inches(0.5), Inches(1.55), Inches(12.3), Inches(5.5),
        font_size=17, color=DARK_GRAY
    )

    # Right-side accent panel
    add_rect(slide, Inches(12.73), Inches(1.3),
             Inches(0.6), Inches(5.9), LIGHT_BLUE)


def slide_04_rq(prs: Presentation):
    """Slide 4 – Research Questions & Objectives."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Research Questions & Objectives")
    add_footer(slide, 4)

    # Primary RQ box
    add_rect(slide, Inches(0.4), Inches(1.5),
             Inches(12.53), Inches(0.85), LIGHT_BLUE)
    add_text_box(
        slide,
        "Primary Research Question:  How does Artificial Intelligence "
        "diffuse spatially across the United Kingdom?",
        Inches(0.6), Inches(1.55), Inches(12.13), Inches(0.75),
        font_size=16, bold=True, color=DARK_BLUE
    )

    # Sub-questions
    sub_q = [
        "Does AI adoption exhibit hierarchical diffusion "
        "(London → secondary cities)?",
        "What spatial concentration patterns emerge at the "
        "UK Outcode level?",
        "Which economic sectors show the highest AI penetration, "
        "and how does this vary by region?",
        "Is there evidence of knowledge spillovers / epidemic "
        "diffusion effects around major clusters?",
    ]
    add_text_box(slide, "Sub-Questions:",
                 Inches(0.5), Inches(2.5), Inches(6.0), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_list(
        slide, sub_q,
        Inches(0.5), Inches(2.95), Inches(6.0), Inches(3.8),
        font_size=14, color=DARK_GRAY
    )

    # Objectives
    objectives = [
        "MAP  – Identify AI firm locations at UK Outcode level",
        "MODEL  – Classify AI vs non-AI firms (ML pipeline)",
        "ANALYSE  – Quantify spatial concentration via LQ",
        "INTERPRET  – Test hierarchical & epidemic diffusion",
        "VISUALISE  – Produce Tableau-ready hotspot dataset",
    ]
    add_text_box(slide, "Objectives:",
                 Inches(6.8), Inches(2.5), Inches(6.0), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_list(
        slide, objectives,
        Inches(6.8), Inches(2.95), Inches(6.0), Inches(3.8),
        font_size=14, color=DARK_GRAY
    )


def slide_05_litreview(prs: Presentation):
    """Slide 5 – Literature Review / Theoretical Framework."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Literature Review – Theoretical Framework",
                  "Foundations from innovation geography & diffusion theory")
    add_footer(slide, 5)

    frameworks = [
        ("Innovation Diffusion Theory",
         "Rogers (1962) & Hägerstrand (1967): innovations spread through "
         "social networks and geographic proximity over time."),
        ("Hierarchical Diffusion",
         "Adoption 'jumps' down the urban hierarchy from primate cities "
         "(London) to secondary centres before spreading to periphery."),
        ("Contagious / Epidemic Diffusion",
         "Geographic clustering via proximity; knowledge spillovers "
         "accelerate local adoption — analogous to disease contagion."),
        ("General Purpose Technology (GPT) Theory",
         "Bresnahan & Trajtenberg (1995): AI exhibits pervasive "
         "applicability, improvement over time, and innovation "
         "complementarities across sectors."),
        ("Urban Hierarchy & Agglomeration Economics",
         "Concentration in dense urban labour markets (Marshall externalities); "
         "talent clustering, thick markets, knowledge exchange."),
    ]

    top = Inches(1.5)
    for title, body in frameworks:
        add_rect(slide, Inches(0.4), top, Inches(12.53), Inches(0.2),
                 DARK_BLUE)
        add_text_box(slide, title,
                     Inches(0.5), top, Inches(12.33), Inches(0.3),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, body,
                     Inches(0.5), top + Inches(0.22), Inches(12.33), Inches(0.55),
                     font_size=13, color=DARK_GRAY)
        top += Inches(0.85)


def slide_06_data(prs: Presentation):
    """Slide 6 – Data & Data Collection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Data & Data Collection",
                  "Firm-level web-scraped dataset of UK AI firms")
    add_footer(slide, 6)

    left_bullets = [
        "Source: Web-scraped UK business websites "
        "(ai_websites.csv)",
        "Key fields: URL, full-text content, postcode, "
        "inferred base domain",
        "National coverage — businesses across all UK "
        "regions and sectors",
        "Post-deduplication: unique firms resolved via "
        "base domain aggregation",
        "Ground-truth labels: AI / Non-AI (human "
        "annotation + keyword heuristics)",
    ]
    add_text_box(slide, "Dataset Overview",
                 Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_list(
        slide, left_bullets,
        Inches(0.5), Inches(1.95), Inches(5.8), Inches(4.8),
        font_size=14, color=DARK_GRAY
    )

    # Right panel – data pipeline steps
    add_rect(slide, Inches(6.7), Inches(1.45),
             Inches(6.2), Inches(5.7), LIGHT_GRAY)
    steps = [
        "① Web Scraping — collect URLs & text",
        "② HTML Parsing — extract clean text",
        "③ Base Domain Extraction — entity resolution",
        "④ Postcode Extraction & Cleaning",
        "⑤ Keyword Heuristics — initial AI labels",
        "⑥ Train / Test Split (80 / 20)",
        "⑦ TF-IDF Features + ML Classification",
    ]
    add_text_box(slide, "Data Pipeline",
                 Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_list(
        slide, steps,
        Inches(6.9), Inches(2.0), Inches(5.8), Inches(4.8),
        font_size=14, color=DARK_GRAY
    )


def slide_07_methodology(prs: Presentation):
    """Slide 7 – Methodology Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Methodology Overview",
                  "Integrated ML + Spatial Econometrics pipeline")
    add_footer(slide, 7)

    stages = [
        ("Data\nCollection",   MID_BLUE),
        ("Entity\nResolution",  DARK_BLUE),
        ("AI\nClassification",  MID_BLUE),
        ("Topic\nModelling",    DARK_BLUE),
        ("Spatial\nAnalysis",   MID_BLUE),
        ("Visualisation",       DARK_BLUE),
    ]

    box_w = Inches(1.8)
    box_h = Inches(1.1)
    top = Inches(2.5)
    gap = Inches(0.2)
    start_left = Inches(0.35)

    for i, (label, color) in enumerate(stages):
        left = start_left + i * (box_w + gap)
        add_rect(slide, left, top, box_w, box_h, color)
        add_text_box(
            slide, label,
            left, top, box_w, box_h,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER
        )
        # Arrow (except after last box)
        if i < len(stages) - 1:
            arr_left = left + box_w
            add_text_box(
                slide, "▶",
                arr_left, top + Inches(0.3), gap, Inches(0.5),
                font_size=16, color=ACCENT_GOLD, align=PP_ALIGN.CENTER
            )

    # Descriptions below pipeline
    descriptions = [
        "Web scraping UK business sites; extract URL, text, postcode",
        "Base domain aggregation; HQ postcode via Mode proxy",
        "TF-IDF + MiniBatchKMeans + Random Forest / Logistic Regression",
        "Latent Dirichlet Allocation; 5 AI sub-sector topics",
        "Outcode aggregation; Location Quotient; pgeocode geocoding",
        "Tableau-ready CSV with LQ values, coordinates & sector labels",
    ]
    desc_top = Inches(3.85)
    for i, desc in enumerate(descriptions):
        left = start_left + i * (box_w + gap)
        add_text_box(
            slide, desc,
            left, desc_top, box_w, Inches(2.0),
            font_size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER
        )

    # Key point
    add_rect(slide, Inches(0.4), Inches(6.4), Inches(12.53), Inches(0.7),
             LIGHT_BLUE)
    add_text_box(
        slide,
        "Key design choice: Outcode (Postal District) chosen as the spatial "
        "unit — granular enough for hotspot detection, yet robust enough to "
        "avoid small-number statistical artefacts.",
        Inches(0.6), Inches(6.42), Inches(12.13), Inches(0.65),
        font_size=13, italic=True, color=DARK_BLUE
    )


def slide_08_entity(prs: Presentation):
    """Slide 8 – Entity Resolution & HQ Proxy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Entity Resolution & HQ Proxy Modelling",
                  "Preventing over-counting; assigning a single location per firm")
    add_footer(slide, 8)

    # Left: problem statement
    add_text_box(slide, "The Challenge",
                 Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    prob = [
        "Multiple sub-pages scraped per firm → naive analysis "
        "over-counts large companies",
        "Different pages may carry different postcodes (branch "
        "offices, service areas)",
        "Need one canonical location per firm for spatial analysis",
    ]
    add_bullet_list(slide, prob,
                    Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.5),
                    font_size=14, color=DARK_GRAY)

    # Left: solution
    add_text_box(slide, "Solution: Base Domain Aggregation",
                 Inches(0.5), Inches(4.6), Inches(5.8), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    sol = [
        "Extract base domain from each URL "
        "(e.g. www.example.co.uk → example.co.uk)",
        "Aggregate all text and postcodes per base domain",
        "Assign HQ postcode = statistical Mode of all postcodes "
        "observed for that domain",
    ]
    add_bullet_list(slide, sol,
                    Inches(0.5), Inches(5.1), Inches(5.8), Inches(2.0),
                    font_size=14, color=DARK_GRAY)

    # Right: example box
    add_rect(slide, Inches(6.8), Inches(1.5),
             Inches(6.1), Inches(5.7), LIGHT_GRAY)
    add_text_box(slide, "Worked Example",
                 Inches(7.0), Inches(1.55), Inches(5.7), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    example_lines = [
        "URL 1: www.deepmind.com/research    → EC1A 1BB",
        "URL 2: www.deepmind.com/careers     → EC1A 1BB",
        "URL 3: www.deepmind.com/blog        → W1A 0AX",
        "URL 4: www.deepmind.com/safety      → EC1A 1BB",
        "",
        "Base domain: deepmind.com",
        "Postcodes observed: [EC1A, EC1A, W1A, EC1A]",
        "Mode → HQ Postcode: EC1A  ✓",
        "",
        "Result: 1 firm · 1 postcode · 1 spatial record",
    ]
    add_bullet_list(slide, example_lines,
                    Inches(7.0), Inches(2.05), Inches(5.7), Inches(4.8),
                    font_size=13, color=DARK_GRAY)


def slide_09_classification(prs: Presentation):
    """Slide 9 – AI Classification Model."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "AI Classification Model",
                  "TF-IDF + MiniBatchKMeans + Ensemble Classifiers  |  95.32% accuracy")
    add_footer(slide, 9)

    # Pipeline steps
    steps = [
        ("TF-IDF Vectorisation",
         "Max 5,000 features; removes stop-words; "
         "represents firm text as numerical feature vectors"),
        ("MiniBatchKMeans Clustering",
         "Unsupervised pre-grouping of text documents "
         "to discover latent groupings; feeds cluster-ID "
         "as additional feature"),
        ("Random Forest Classifier",
         "Ensemble of 100 decision trees; handles class "
         "imbalance; provides feature importances"),
        ("Logistic Regression",
         "Linear baseline; fast inference; "
         "probability calibration for threshold tuning"),
    ]

    top = Inches(1.55)
    colors = [MID_BLUE, DARK_BLUE, MID_BLUE, DARK_BLUE]
    for (title, body), color in zip(steps, colors):
        add_rect(slide, Inches(0.4), top, Inches(0.15), Inches(0.7), color)
        add_text_box(slide, title,
                     Inches(0.65), top, Inches(6.2), Inches(0.32),
                     font_size=14, bold=True, color=color)
        add_text_box(slide, body,
                     Inches(0.65), top + Inches(0.3), Inches(6.2), Inches(0.45),
                     font_size=13, color=DARK_GRAY)
        top += Inches(0.85)

    # Metrics table (right panel)
    add_rect(slide, Inches(7.2), Inches(1.5),
             Inches(5.7), Inches(5.7), LIGHT_GRAY)
    add_text_box(slide, "Model Performance Metrics",
                 Inches(7.4), Inches(1.55), Inches(5.3), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)

    # Table headers
    cols = ["Class", "Precision", "Recall", "F1-Score", "Support"]
    col_widths = [Inches(1.1), Inches(1.1), Inches(1.0), Inches(1.1), Inches(1.1)]
    rows_data = [
        ["Non-AI", "99.6%", "95.6%", "97.6%", "4,892"],
        ["AI",     "29.0%", "80.6%", "42.6%",   "108"],
        ["Accuracy (Overall)", "",    "",    "95.32%",   "5,000"],
    ]

    header_top = Inches(2.05)
    add_rect(slide, Inches(7.35), header_top,
             Inches(5.3), Inches(0.35), DARK_BLUE)
    x = Inches(7.35)
    for i, (col, w) in enumerate(zip(cols, col_widths)):
        add_text_box(slide, col, x, header_top, w, Inches(0.35),
                     font_size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        x += w

    row_colors = [WHITE, LIGHT_GRAY, LIGHT_BLUE]
    for ri, (row, bg) in enumerate(zip(rows_data, row_colors)):
        row_top = header_top + Inches(0.35) + ri * Inches(0.4)
        add_rect(slide, Inches(7.35), row_top, Inches(5.3), Inches(0.4), bg)
        x = Inches(7.35)
        for ci, (cell, w) in enumerate(zip(row, col_widths)):
            cell_color = DARK_GRAY
            if ri == 0 and ci == 1:  # Non-AI precision – highlight green
                cell_color = GREEN
            if ri == 1 and ci == 2:  # AI recall – highlight gold
                cell_color = ACCENT_GOLD
            add_text_box(slide, cell, x, row_top, w, Inches(0.4),
                         font_size=12, bold=(ri == 2),
                         color=cell_color, align=PP_ALIGN.CENTER)
            x += w

    # Note on class imbalance
    add_text_box(
        slide,
        "Note: Class imbalance (108 AI vs 4,892 Non-AI) explains the low "
        "AI precision. High AI recall (80.6%) is prioritised to minimise "
        "missed true AI firms.",
        Inches(7.35), Inches(3.7), Inches(5.3), Inches(1.0),
        font_size=11, italic=True, color=MID_GRAY
    )


def slide_10_lda(prs: Presentation):
    """Slide 10 – Topic Modelling – LDA Results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Topic Modelling – LDA Results",
                  "Latent Dirichlet Allocation reveals 5 AI sub-sectors")
    add_footer(slide, 10)

    topics = [
        ("Topic 1 – Management",
         MID_BLUE,
         "business, strategy, operations, consultancy, talent, leadership, "
         "decision-making, enterprise, governance, analytics"),
        ("Topic 2 – Commerce",
         ACCENT_GOLD,
         "retail, e-commerce, customer, product, sales, marketing, "
         "recommendation, personalisation, fraud, payment"),
        ("Topic 3 – Engineering",
         GREEN,
         "software, hardware, robotics, sensor, simulation, design, "
         "systems, manufacturing, control, autonomous"),
        ("Topic 4 – Infrastructure",
         DARK_BLUE,
         "cloud, data, platform, API, server, network, security, "
         "compute, storage, deployment, MLOps"),
        ("Topic 5 – Automation",
         RED,
         "process, workflow, RPA, document, extraction, pipeline, "
         "scheduling, rule, trigger, efficiency"),
    ]

    for i, (title, color, keywords) in enumerate(topics):
        col = i % 3
        row = i // 3
        left = Inches(0.3 + col * 4.35)
        top  = Inches(1.6 + row * 2.35)

        add_rect(slide, left, top, Inches(4.1), Inches(0.35), color)
        add_text_box(slide, title, left + Inches(0.1), top,
                     Inches(3.9), Inches(0.35),
                     font_size=13, bold=True, color=WHITE)
        add_rect(slide, left, top + Inches(0.35), Inches(4.1), Inches(1.9),
                 LIGHT_GRAY)
        add_text_box(
            slide,
            f"Top keywords:\n{keywords}",
            left + Inches(0.1), top + Inches(0.4),
            Inches(3.9), Inches(1.8),
            font_size=11, color=DARK_GRAY
        )

    # GPT interpretation
    add_rect(slide, Inches(0.3), Inches(6.45),
             Inches(12.73), Inches(0.65), LIGHT_BLUE)
    add_text_box(
        slide,
        "Interpretation: Five distinct sub-sectors spanning Management, Commerce, "
        "Engineering, Infrastructure, and Automation confirm AI's role as a "
        "General Purpose Technology (GPT) — pervasive across the entire economy.",
        Inches(0.5), Inches(6.48), Inches(12.33), Inches(0.6),
        font_size=13, bold=True, color=DARK_BLUE
    )


def slide_11_spatial_method(prs: Presentation):
    """Slide 11 – Spatial Analysis Methodology."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Spatial Analysis Methodology",
                  "From firm-level data to UK AI hotspot maps")
    add_footer(slide, 11)

    # Left: steps
    add_text_box(slide, "Analytical Steps",
                 Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)
    steps = [
        "Aggregate classified AI firms to UK Outcode "
        "(e.g. EC1, SW1, M1) — 2,900+ districts",
        "Calculate total firm count per Outcode as "
        "denominator (all scraped businesses)",
        "Compute Location Quotient (LQ) per Outcode",
        "Geocode each Outcode centroid using pgeocode",
        "Filter AI hotspots: retain Outcodes with LQ ≥ 1.0",
        "Export Tableau-ready CSV with lat/lon, LQ, "
        "sector composition",
    ]
    add_bullet_list(slide, steps,
                    Inches(0.5), Inches(2.05), Inches(5.8), Inches(4.8),
                    font_size=14, color=DARK_GRAY)

    # Right: LQ formula
    add_rect(slide, Inches(6.8), Inches(1.5),
             Inches(6.1), Inches(5.7), LIGHT_BLUE)
    add_text_box(slide, "Location Quotient (LQ) Formula",
                 Inches(7.0), Inches(1.6), Inches(5.7), Inches(0.4),
                 font_size=15, bold=True, color=DARK_BLUE)

    add_text_box(
        slide,
        "LQ  =  (AI firms in Outcode / Total firms in Outcode)\n"
        "        ÷\n"
        "       (Total AI firms in UK / Total firms in UK)",
        Inches(7.0), Inches(2.15), Inches(5.7), Inches(1.1),
        font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER
    )

    lq_interp = [
        "LQ > 1.0  →  Outcode has above-average AI concentration  ✓",
        "LQ = 1.0  →  National average",
        "LQ < 1.0  →  Below-average AI concentration",
        "LQ ≥ 2.0  →  Significant AI cluster (hotspot)",
    ]
    add_bullet_list(slide, lq_interp,
                    Inches(7.0), Inches(3.4), Inches(5.7), Inches(2.5),
                    font_size=13, color=DARK_GRAY)

    add_text_box(
        slide,
        "Spatial unit: Outcode chosen over LAD or LSOA for optimal "
        "balance between granularity and statistical robustness.",
        Inches(7.0), Inches(6.0), Inches(5.7), Inches(0.8),
        font_size=12, italic=True, color=MID_GRAY
    )


def slide_12_findings_spatial(prs: Presentation):
    """Slide 12 – Key Findings: Spatial Patterns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Key Findings – Spatial Patterns",
                  "AI hotspot geography across the UK")
    add_footer(slide, 12)

    findings = [
        ("London Dominance",
         "London Outcodes (EC, WC, W, SW, SE, N, E) account for the "
         "majority of AI hotspots by LQ and absolute count. The "
         "City of London and Tech City / Silicon Roundabout "
         "emerge as primary clusters."),
        ("Hierarchical Diffusion Evidence",
         "Secondary cities — Manchester (M1–M4), Edinburgh (EH), "
         "Bristol (BS), Leeds (LS) — show elevated LQs, consistent "
         "with hierarchical trickle-down from London."),
        ("Geographic Clustering",
         "Spatial autocorrelation evident: AI hotspot Outcodes "
         "tend to be adjacent to other hotspot Outcodes, "
         "suggesting localised knowledge spillovers."),
        ("National AI Density",
         "Hotspot dataset (LQ ≥ 1.0) captures significant share "
         "of classified AI firms. Peripheral UK regions show "
         "markedly lower LQ values — a digital divide."),
        ("Output: ai_hotspots_for_tableau.csv",
         "Final dataset with Outcode, LQ, firm count, lat/lon, "
         "and sector composition — ready for Tableau choropleth "
         "and dot-density mapping."),
    ]

    top = Inches(1.55)
    alt = True
    for title, body in findings:
        bg = LIGHT_GRAY if alt else WHITE
        add_rect(slide, Inches(0.4), top, Inches(12.53), Inches(0.9), bg)
        add_text_box(slide, title,
                     Inches(0.6), top + Inches(0.05), Inches(3.5), Inches(0.35),
                     font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, body,
                     Inches(4.0), top + Inches(0.05), Inches(8.7), Inches(0.8),
                     font_size=13, color=DARK_GRAY)
        top += Inches(0.95)
        alt = not alt


def slide_13_findings_sectoral(prs: Presentation):
    """Slide 13 – Key Findings: Sectoral Analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Key Findings – Sectoral Analysis",
                  "LDA sub-sectors reveal AI's role as a General Purpose Technology")
    add_footer(slide, 13)

    # Five sector boxes
    sectors = [
        ("Management",    "37%", MID_BLUE,
         "Widely distributed nationally; strong presence in financial "
         "districts and corporate headquarters."),
        ("Commerce",      "24%", ACCENT_GOLD,
         "Concentrated in retail / logistics corridors; "
         "Manchester and Birmingham show high relative density."),
        ("Engineering",   "15%", GREEN,
         "Clusters near university cities (Cambridge, Bristol, "
         "Edinburgh) and manufacturing belts."),
        ("Infrastructure","14%", DARK_BLUE,
         "Dominant in London data-centre corridors; "
         "cloud and MLOps firms in EC and W postcodes."),
        ("Automation",     "10%", RED,
         "Mixed geography; manufacturing regions (West Midlands, "
         "Sheffield) show elevated representation."),
    ]

    col_w = Inches(2.45)
    for i, (name, pct, color, desc) in enumerate(sectors):
        left = Inches(0.3 + i * 2.6)
        add_rect(slide, left, Inches(1.55), col_w, Inches(0.5), color)
        add_text_box(slide, name, left, Inches(1.55), col_w, Inches(0.5),
                     font_size=14, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_rect(slide, left, Inches(2.05), col_w, Inches(0.55), LIGHT_BLUE)
        add_text_box(slide, pct, left, Inches(2.05), col_w, Inches(0.55),
                     font_size=28, bold=True, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER)
        add_rect(slide, left, Inches(2.6), col_w, Inches(2.5), LIGHT_GRAY)
        add_text_box(slide, desc, left + Inches(0.1), Inches(2.65),
                     col_w - Inches(0.2), Inches(2.4),
                     font_size=12, color=DARK_GRAY)

    # Summary
    add_rect(slide, Inches(0.3), Inches(5.3), Inches(12.73), Inches(0.9),
             LIGHT_BLUE)
    add_text_box(
        slide,
        "GPT Interpretation: The presence of AI across five economically "
        "distinct sub-sectors — from Engineering to Commerce to Management — "
        "empirically supports Bresnahan & Trajtenberg's (1995) GPT framework. "
        "AI is not a niche technology; it is economically pervasive.",
        Inches(0.5), Inches(5.35), Inches(12.33), Inches(0.8),
        font_size=13, italic=True, color=DARK_BLUE
    )

    # Regional note
    add_text_box(
        slide,
        "Regional LQ decomposition by sector enables identification of "
        "sectoral specialisation — e.g., Infrastructure over-represented "
        "in London; Engineering over-represented in Cambridge and Bristol.",
        Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.6),
        font_size=12, color=MID_GRAY
    )


def slide_14_findings_diffusion(prs: Presentation):
    """Slide 14 – Key Findings: Diffusion Patterns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Key Findings – Diffusion Patterns",
                  "Hierarchical & epidemic diffusion both confirmed in the UK AI economy")
    add_footer(slide, 14)

    # Two side-by-side boxes
    for col, (title, color, points) in enumerate([
        ("Hierarchical Diffusion  ✓", DARK_BLUE, [
            "London acts as the unequivocal primary hub (primate city)",
            "AI presence cascades to Tier-2 cities: Manchester, "
            "Edinburgh, Bristol, Leeds, Birmingham",
            "Consistent with Hägerstrand's hierarchical wave model",
            "Urban rank closely predicts AI firm density",
            "Implies policy levers must target secondary cities "
            "to achieve equitable diffusion",
        ]),
        ("Epidemic / Contagious Diffusion  ✓", MID_BLUE, [
            "Spatial clustering of hotspot Outcodes — "
            "adjacency drives adoption",
            "Knowledge spillovers evident around Tech City "
            "(EC1 / N1) and Cambridge (CB)",
            "Proximity to existing AI firms significantly "
            "predicts new firm emergence",
            "Consistent with epidemic model: 'infection' "
            "spreads through geographic contact",
            "Agglomeration externalities (talent, capital, "
            "supply chains) reinforce clustering",
        ]),
    ]):
        left = Inches(0.3 + col * 6.5)
        add_rect(slide, left, Inches(1.55),
                 Inches(6.2), Inches(0.45), color)
        add_text_box(slide, title,
                     left + Inches(0.1), Inches(1.55),
                     Inches(6.0), Inches(0.45),
                     font_size=14, bold=True, color=WHITE)
        add_bullet_list(
            slide, points,
            left + Inches(0.1), Inches(2.1),
            Inches(6.0), Inches(4.5),
            font_size=13, color=DARK_GRAY
        )

    # GPT footnote
    add_rect(slide, Inches(0.3), Inches(6.45),
             Inches(12.73), Inches(0.65), ACCENT_GOLD)
    add_text_box(
        slide,
        "Both diffusion mechanisms operate simultaneously — hierarchical "
        "at the macro (city) scale, epidemic at the micro (neighbourhood) "
        "scale — mirroring findings in prior innovation geography literature.",
        Inches(0.5), Inches(6.48), Inches(12.33), Inches(0.6),
        font_size=13, bold=True, color=DARK_BLUE
    )


def slide_15_discussion(prs: Presentation):
    """Slide 15 – Discussion & Implications."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Discussion & Implications",
                  "Theoretical contributions and policy relevance")
    add_footer(slide, 15)

    sections = [
        ("Theoretical Contributions", DARK_BLUE, [
            "First large-scale firm-level spatial analysis of UK AI diffusion",
            "Integrates ML classification with spatial econometrics in a "
            "single reproducible pipeline",
            "Empirically validates both hierarchical and epidemic diffusion "
            "models in the AI context",
            "Confirms GPT pervasiveness via LDA sub-sector analysis",
        ]),
        ("Policy Implications", MID_BLUE, [
            "Regional AI investment should target secondary cities to "
            "stimulate hierarchical trickle-down",
            "Co-location incentives (innovation districts, enterprise zones) "
            "can amplify epidemic spillovers",
            "Digital divide risk: peripheral regions require targeted "
            "skills and infrastructure investment",
            "UK AI Strategy should disaggregate national targets to "
            "Outcode / LAD level for spatial accountability",
        ]),
        ("Broader Significance", ACCENT_GOLD, [
            "Framework is sector-agnostic and can be applied to other "
            "emerging technologies (quantum, biotech)",
            "Postcode-level granularity enables precise local authority "
            "decision-making",
            "Open-source pipeline (Python) supports reproducibility "
            "and future longitudinal extensions",
        ]),
    ]

    top = Inches(1.55)
    for title, color, points in sections:
        add_rect(slide, Inches(0.4), top, Inches(12.53), Inches(0.38), color)
        add_text_box(slide, title,
                     Inches(0.6), top, Inches(12.13), Inches(0.38),
                     font_size=14, bold=True, color=WHITE)
        add_bullet_list(
            slide, points,
            Inches(0.6), top + Inches(0.38), Inches(12.13),
            Inches(len(points) * 0.35 + 0.1),
            font_size=13, color=DARK_GRAY
        )
        top += Inches(0.38 + len(points) * 0.35 + 0.2)


def slide_16_limitations(prs: Presentation):
    """Slide 16 – Limitations & Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Limitations & Future Work",
                  "Honest appraisal and roadmap for future research")
    add_footer(slide, 16)

    limitations = [
        ("Web Scraping Coverage",
         "Dataset limited to publicly accessible business websites. "
         "Stealth AI R&D (in-house corporate AI, government) "
         "likely under-represented."),
        ("Temporal Snapshot",
         "Data reflects a single point in time. Dynamic diffusion "
         "patterns (velocity, acceleration) cannot be inferred without "
         "longitudinal panel data."),
        ("HQ Proxy Accuracy",
         "Mode-based postcode proxy may misassign HQ for firms with "
         "multiple equally-represented locations or remote-first workforces."),
        ("Class Imbalance",
         "AI firms are rare relative to the total business population. "
         "Low AI precision (29%) risks false positives; threshold "
         "calibration is an ongoing challenge."),
        ("Spatial Unit",
         "Outcode aggregation masks within-Outcode heterogeneity. "
         "Smaller units (e.g., unit postcode) would improve resolution "
         "but reduce statistical robustness."),
    ]

    future = [
        "Panel data: track AI firm emergence year-on-year (2015–2025)",
        "Network analysis: map knowledge flows between AI firms",
        "International comparison: UK vs. EU vs. US spatial patterns",
        "Causal inference: DiD analysis of policy interventions",
        "Real-time pipeline: live web-scraping + automated LQ updates",
    ]

    add_text_box(slide, "Limitations",
                 Inches(0.5), Inches(1.5), Inches(6.5), Inches(0.35),
                 font_size=15, bold=True, color=DARK_BLUE)

    top = Inches(1.95)
    for title, body in limitations:
        add_text_box(slide, f"▪  {title}:",
                     Inches(0.5), top, Inches(6.3), Inches(0.3),
                     font_size=13, bold=True, color=DARK_BLUE)
        add_text_box(slide, body,
                     Inches(0.7), top + Inches(0.28), Inches(6.1), Inches(0.5),
                     font_size=12, color=DARK_GRAY)
        top += Inches(0.88)

    add_text_box(slide, "Future Work",
                 Inches(7.2), Inches(1.5), Inches(5.7), Inches(0.35),
                 font_size=15, bold=True, color=DARK_BLUE)
    add_bullet_list(
        slide, future,
        Inches(7.2), Inches(1.95), Inches(5.7), Inches(4.0),
        font_size=14, color=DARK_GRAY
    )


def slide_17_conclusion(prs: Presentation):
    """Slide 17 – Conclusion."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, WHITE)
    add_title_bar(slide, "Conclusion",
                  "Summary of contributions and key takeaways")
    add_footer(slide, 17)

    contributions = [
        ("Contribution 1: Novel Spatial Dataset",
         "First firm-level UK AI hotspot map at Outcode resolution, "
         "produced via an integrated ML + spatial econometrics pipeline."),
        ("Contribution 2: Validated Diffusion Theory",
         "Empirical evidence for both hierarchical and epidemic AI diffusion "
         "in the UK — answering the primary research question."),
        ("Contribution 3: GPT Confirmation",
         "LDA topic modelling across 5 distinct sub-sectors provides "
         "firm-level evidence of AI's General Purpose Technology status."),
        ("Contribution 4: Reproducible Open Pipeline",
         "Python codebase (Model.py, Topic_Modeling.py, Spatial_Analysis.py) "
         "is modular, scalable, and applicable to other technology sectors."),
        ("Contribution 5: Policy-Ready Output",
         "Tableau-ready CSV (ai_hotspots_for_tableau.csv) provides actionable "
         "spatial intelligence for UK regional AI strategy."),
    ]

    top = Inches(1.6)
    for i, (title, body) in enumerate(contributions):
        bg = LIGHT_BLUE if i % 2 == 0 else LIGHT_GRAY
        add_rect(slide, Inches(0.4), top, Inches(12.53), Inches(0.85), bg)
        add_text_box(slide, title,
                     Inches(0.6), top + Inches(0.05), Inches(12.13), Inches(0.32),
                     font_size=14, bold=True, color=DARK_BLUE)
        add_text_box(slide, body,
                     Inches(0.6), top + Inches(0.37), Inches(12.13), Inches(0.45),
                     font_size=13, color=DARK_GRAY)
        top += Inches(0.92)

    # Closing statement
    add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.53), Inches(0.65),
             DARK_BLUE)
    add_text_box(
        slide,
        "\"AI diffusion in the UK is spatially uneven, hierarchically "
        "structured, and sectorally pervasive — understanding this geography "
        "is essential for equitable and effective AI policy.\"",
        Inches(0.6), Inches(6.52), Inches(12.13), Inches(0.6),
        font_size=13, bold=True, italic=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )


def slide_18_references(prs: Presentation):
    """Slide 18 – References & Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg_color(slide, DARK_BLUE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.1), ACCENT_GOLD)
    add_rect(slide, Inches(0), Inches(7.4), SLIDE_W, Inches(0.1), ACCENT_GOLD)

    add_text_box(
        slide, "References & Acknowledgements",
        Inches(0.5), Inches(0.25), Inches(12.33), Inches(0.55),
        font_size=26, bold=True, color=WHITE
    )

    refs = [
        "Bresnahan, T.F. & Trajtenberg, M. (1995). General purpose technologies: "
        "'Engines of growth'? Journal of Econometrics, 65(1), 83–108.",
        "Rogers, E.M. (1962). Diffusion of Innovations. Free Press, New York.",
        "Hägerstrand, T. (1967). Innovation Diffusion as a Spatial Process. "
        "University of Chicago Press.",
        "Audretsch, D.B. & Feldman, M.P. (1996). R&D spillovers and the geography "
        "of innovation and production. American Economic Review, 86(3), 630–640.",
        "Bathelt, H., Malmberg, A. & Maskell, P. (2004). Clusters and knowledge: "
        "local buzz, global pipelines and the process of knowledge creation. "
        "Progress in Human Geography, 28(1), 31–56.",
        "Blei, D.M., Ng, A.Y. & Jordan, M.I. (2003). Latent Dirichlet Allocation. "
        "Journal of Machine Learning Research, 3, 993–1022.",
        "Dosi, G. (1982). Technological paradigms and technological trajectories. "
        "Research Policy, 11(3), 147–162.",
        "DCMS / DSIT (2023). AI Activity in UK Businesses. "
        "UK Department for Science, Innovation & Technology.",
    ]

    add_bullet_list(
        slide, refs,
        Inches(0.5), Inches(0.9), Inches(12.33), Inches(5.0),
        font_size=11, color=LIGHT_BLUE
    )

    # Thank you panel
    add_rect(slide, Inches(2.0), Inches(5.8), Inches(9.33), Inches(1.4),
             WHITE)
    add_text_box(
        slide, "Thank You",
        Inches(2.0), Inches(5.85), Inches(9.33), Inches(0.65),
        font_size=36, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER
    )
    add_text_box(
        slide, "Questions & Discussion Welcome",
        Inches(2.0), Inches(6.55), Inches(9.33), Inches(0.5),
        font_size=18, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER
    )


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def build_presentation(output_path: str = "dissertation_presentation.pptx"):
    """Build the full presentation and save it to output_path."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_01_title(prs)
    print("  [1/18] Title slide")
    slide_02_toc(prs)
    print("  [2/18] Table of Contents")
    slide_03_context(prs)
    print("  [3/18] Research Context & Motivation")
    slide_04_rq(prs)
    print("  [4/18] Research Questions & Objectives")
    slide_05_litreview(prs)
    print("  [5/18] Literature Review")
    slide_06_data(prs)
    print("  [6/18] Data & Data Collection")
    slide_07_methodology(prs)
    print("  [7/18] Methodology Overview")
    slide_08_entity(prs)
    print("  [8/18] Entity Resolution & HQ Proxy")
    slide_09_classification(prs)
    print("  [9/18] AI Classification Model")
    slide_10_lda(prs)
    print(" [10/18] Topic Modelling – LDA")
    slide_11_spatial_method(prs)
    print(" [11/18] Spatial Analysis Methodology")
    slide_12_findings_spatial(prs)
    print(" [12/18] Findings – Spatial Patterns")
    slide_13_findings_sectoral(prs)
    print(" [13/18] Findings – Sectoral Analysis")
    slide_14_findings_diffusion(prs)
    print(" [14/18] Findings – Diffusion Patterns")
    slide_15_discussion(prs)
    print(" [15/18] Discussion & Implications")
    slide_16_limitations(prs)
    print(" [16/18] Limitations & Future Work")
    slide_17_conclusion(prs)
    print(" [17/18] Conclusion")
    slide_18_references(prs)
    print(" [18/18] References & Thank You")

    prs.save(output_path)
    print(f"\nPresentation saved → {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
